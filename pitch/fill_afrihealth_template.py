"""
Fill the white Afrihealth Innovation template with Aurora Travels content.

- Keeps template styling, structure, and decorative images
- Copies details from the filled AuroraTravels pitch
- Forces all text bold + dark/legible ink
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent
SRC = Path(
    "/home/ubuntu/.cursor/projects/workspace/uploads/"
    "Afrihealth_Innovation_Template_-_MAKE_A_COPY_3605.pptx"
)
OUT = ROOT / "AuroraTravels_Pitch_Deck.pptx"
TEMPLATE_COPY = ROOT / "Afrihealth_Innovation_Template.pptx"

# Near-black ink for legibility on white slides
INK = RGBColor(0x14, 0x14, 0x14)
MIN_BODY_PT = 11


def set_text(shape, text: str) -> None:
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    lines = text.split("\n")
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()
    for i, para in enumerate(tf.paragraphs):
        if i < len(lines):
            line = lines[i]
            if para.runs:
                para.runs[0].text = line
                for r in para.runs[1:]:
                    r.text = ""
            else:
                para.add_run().text = line
        else:
            for r in para.runs:
                r.text = ""


def first_contains(slide, fragment: str):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs)
        if fragment in full:
            return shape
    return None


def first_exact(slide, exact: str):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
        if full == exact.strip():
            return shape
    return None


def make_legible(prs: Presentation) -> None:
    """Force every text run bold + dark ink; bump tiny body sizes."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = True
                    # Solid dark ink (override washed-out greys)
                    run.font.color.rgb = INK
                    try:
                        size = run.font.size
                        if size is not None and size < Pt(MIN_BODY_PT):
                            # Keep tiny labels from going unreadably small
                            if size < Pt(9):
                                run.font.size = Pt(10)
                    except Exception:
                        pass


def fill(prs: Presentation) -> None:
    # ---- 1 Title ----
    s = prs.slides[0]
    set_text(first_contains(s, "PITCH DECK TEMPLATE"), "AURORA TRAVELS")
    set_text(
        first_contains(s, "Afrihealth Innovation Challenge"),
        "Afrihealth Innovation Challenge 2026",
    )
    set_text(
        first_contains(s, "Replace the guidance text"),
        "KENYA · TRAVEL · CULTURE · COMMERCE\n"
        "A single, Kenya-native app for parks, stays, verified crafts and student guides — "
        "pitched for partnership and investment.\n"
        "Discover 15 parks · Shop 28 verified crafts · Pay with native M-Pesa STK Push · "
        "Connect 12 university guides\n"
        "Built in Nairobi by ProximaOpal — web & AI product engineer · auroratravels",
    )

    # ---- 2 Problem ----
    s = prs.slides[1]
    set_text(
        first_contains(s, "What health challenge"),
        "Kenya's travel experience is fragmented — and its creative economy is invisible online",
    )
    set_text(
        first_contains(s, "Describe the specific health"),
        "• International arrivals rose 9% year-on-year to 2.7 million in 2025 (from 2.47M in 2024), "
        "with 5.2M domestic trips — 7.9M total visitors — yet travelers still stitch together "
        "separate tools for parks, stays, curio shopping and guides, mostly through "
        "commission-heavy foreign booking platforms.\n"
        "• Kenya's informal / Jua Kali and MSE economy contributes an estimated 24–33% of GDP "
        "and the bulk of new jobs, but most craft artisans still lack a verified, direct digital "
        "storefront into tourist demand.\n"
        "• Federation of Kenya Employers cites a 67% unemployment rate among Kenyans aged 15–34 — "
        "deep local knowledge sits idle instead of being monetized as guiding services.",
    )
    set_text(
        first_contains(s, "Supporting data point"),
        "Supporting data point\n"
        "67% — FKE figure for Kenyans aged 15–34 without adequate employment\n"
        "2.7M international arrivals in 2025, up 9% YoY; KES 0.5T tourism earnings\n"
        "Sources: Kenya Ministry of Tourism & Wildlife / Magical Kenya 2025; FKE; FinAccess / "
        "Daily Nation on informal GDP share (24–33%).",
    )

    # ---- 3 Solution ----
    s = prs.slides[2]
    set_text(
        first_contains(s, "How does your innovation"),
        "One app: parks, stays, verified crafts and student guides — native to Kenya",
    )
    set_text(
        first_contains(s, "One-sentence description"),
        "A single web experience from discovering Kenya's parks to buying a named artisan craft "
        "and pairing with a vetted university guide — checked out with native M-Pesa STK Push.",
    )
    set_text(
        first_contains(s, "The core mechanism"),
        "Interactive map across 15 parks → unified travel & stay booking → curated marketplace of "
        "28 crafts tied to real shops and Google Maps pins → on-demand pairing with vetted "
        "university student guides.",
    )
    set_text(
        first_contains(s, "The unique insight"),
        "Built on Kenya's own payment rail (mobile money used by ~98% of mobile subscriptions, "
        "CA Kenya Dec 2025) instead of card-only global platforms; every shop and guide is named "
        "and mappable; guiding fees go to Kenyan students.",
    )

    # ---- 4 Product ----
    s = prs.slides[3]
    set_text(
        first_contains(s, "Show the user journey"),
        "Four steps, one continuous journey",
    )
    set_text(
        first_contains(s, "Insert 3–5 step"),
        "The live build today: an interactive map, integrated booking, a working M-Pesa checkout, "
        "and a guide roster — with Kenyan craft & place photography throughout.",
    )
    step_titles, step_descs = [], []
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        t = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
        if t.startswith("Step ") and len(t) <= 8:
            step_titles.append(shape)
        elif t == "[brief description]":
            step_descs.append(shape)
    for i, label in enumerate(["Discover", "Plan", "Shop", "Connect"]):
        if i < len(step_titles):
            set_text(step_titles[i], label)
    bodies = [
        "Browse 15 national parks & destinations across Kenya and Africa on a live interactive map, with search and weather.",
        "Book travel and accommodation in the same flow — no separate app or tab.",
        "Buy an authentic craft from a named shop, e.g. Tabaka Soapstone Cooperative in Kisii, mapped on Google Maps.",
        "Pay instantly by M-Pesa STK Push, then pair with a vetted student guide from one of 12 Kenyan universities.",
    ]
    for i, body in enumerate(bodies):
        if i < len(step_descs):
            set_text(step_descs[i], body)

    # ---- 5 Market ----
    s = prs.slides[4]
    set_text(
        first_contains(s, "Replace with your figures"),
        "TAM — US$12.7B   ·   SAM — ~US$3.8B (KES 0.5T)   ·   SOM (3-yr, modeled) — US$45M GMV",
    )
    set_text(
        first_contains(s, "Total Addressable Market"),
        "• TAM — US$12.7B: WTTC reports Travel & Tourism contributed US$12.7 billion to Kenya's "
        "economy in 2025 (9.3% of GDP) and supported 1.8 million jobs (8.3% of employment). "
        "Source: WTTC Economic Impact Research 2025/26.\n"
        "• SAM — ~US$3.8B (KES 0.5T): Government-reported direct tourism earnings from 7.9M visitors "
        "in 2025 (2.7M international + 5.2M domestic). Digitally reachable craft, guiding and "
        "local-booking spend sits inside this earnings base. Source: Kenya Ministry of Tourism & "
        "Wildlife / Magical Kenya 2025.\n"
        "• SOM (3-yr, modeled) — US$45M GMV: Illustrative ~1.2% of international visitor spend on "
        "crafts/experiences/guiding (~US$5.0B international spend per WTTC) flowing through Aurora "
        "by year 3 — a model to validate in the Nairobi + Maasai Mara pilot, not an achieved figure.",
    )

    # ---- 6 Business model ----
    s = prs.slides[5]
    set_text(first_contains(s, "How do you make money"), "How Aurora Travels makes money")
    set_text(
        first_contains(s, "Subscription, transaction fee"),
        "12% average take-rate on craft GMV (below typical OTA bands); KES 500–800 guide-booking "
        "service fee per pairing; 8–10% referral commission on travel & stay bookings; "
        "featured-placement fees for shops (KES 3,000–8,000/mo) and guides.",
    )
    set_text(
        first_contains(s, "What do customers pay"),
        "Live catalogue prices run from KES 2,200 to KES 15,200 (~US$17–$118) across 28 crafts. "
        "Global tour OTAs commonly take 20–30% (Viator ~20–25% base; GetYourGuide 20–30%). "
        "Aurora undercuts that with M-Pesa-native checkout — no card gateway friction where "
        "~98% of mobile subscriptions use mobile money.",
    )
    set_text(
        first_contains(s, "Customer acquisition cost"),
        "Low CAC via university guide network and shops already mapped. At modeled Y3 GMV of "
        "US$45M and 11% blended take-rate → ~US$5.0M platform revenue. LTV compounds across "
        "repeat trips, multiple crafts and guide bookings.",
    )

    # ---- 7 Traction ----
    s = prs.slides[6]
    set_text(
        first_contains(s, "Proof that this is working"),
        "What's already built — pre-launch, seeking pilot partners",
    )
    nums = []
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        t = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
        if t == "[ number ]":
            nums.append(shape)
    for i, val in enumerate(["28", "15", "12", "1"]):
        if i < len(nums):
            set_text(nums[i], val)
    for old, new in [
        (
            "Users / Patients Reached",
            "Curated crafts listed across 24+ named Kenyan shops (KES 2,200–15,200)",
        ),
        (
            "Pilot Partners",
            "Parks & destinations mapped — 10 in Kenya + 5 across Africa",
        ),
        (
            "Revenue / Funding Raised",
            "Vetted student guides from 12 Kenyan universities",
        ),
        (
            "Key Milestone",
            "Working M-Pesa STK Push — live gateway, phone validation & status polling",
        ),
    ]:
        set_text(first_exact(s, old), new)

    # ---- 8 Competition ----
    s = prs.slides[7]
    set_text(
        first_contains(s, "Axis 2"),
        "Kenya-native payments & verified local shops/guides →",
    )
    set_text(first_contains(s, "Axis 1"), "Breadth of Kenya travel experience ↑")
    set_text(first_exact(s, "You"), "Aurora Travels")
    set_text(first_exact(s, "Comp. A"), "SafariBookings")
    set_text(first_exact(s, "Comp. B"), "Viator / GYG / Klook")
    set_text(first_exact(s, "Comp. C"), "Informal curios")
    set_text(
        first_contains(s, "Label both axes"),
        "• Global OTAs charge 20–30% commission and don't touch M-Pesa or named artisan shops. "
        "Aurora does both, natively — where mobile money reaches ~98% of mobile subscriptions.\n"
        "• SafariBookings aggregates operator tours but does not process payment or sell "
        "crafts/guides directly. Viator, GetYourGuide and Klook are card-first and generic.\n"
        "• Defensible advantage: M-Pesa-native checkout + verified local shops/guides in one "
        "Kenya-first journey.",
    )

    # ---- 9 GTM ----
    s = prs.slides[8]
    set_text(
        first_contains(s, "How you'll acquire"),
        "How Aurora Travels acquires and scales users",
    )
    set_text(
        first_contains(s, "Initial geography"),
        "Nairobi and the Maasai Mara gateway (Narok). Onboard shop partners already mapped — "
        "Maasai Market Nairobi, Tabaka Soapstone (Kisii), Wamunyu Carvers (Machakos), Kazuri Beads — "
        "and the first guide cohort from partner universities. Target: 40 shops, 40 guides, "
        "first 5,000 paid transactions.",
    )
    set_text(
        first_contains(s, "Expansion plan"),
        "Extend across more of Kenya's 47 counties toward the government's 2027 ambition of "
        "5.5M international arrivals and KES 1T tourism earnings. Prioritise 2025 top source "
        "markets: United States, Uganda, Tanzania and the United Kingdom (Ministry report).",
    )
    set_text(
        first_contains(s, "Path to regional"),
        "Expand into the shared East African safari circuit (Tanzania, Uganda) that Aurora's park "
        "map already spans — the Great Rift / Mara–Serengeti ecosystem — with cross-border "
        "M-Pesa and craft corridors.",
    )

    # ---- 10 Team ----
    s = prs.slides[9]
    set_text(first_contains(s, "Why you're the right"), "Why this team can build it")
    photo_ph, name_ph, role_ph, bio_ph = [], [], [], []
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        t = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
        if t == "[ photo ]":
            photo_ph.append(shape)
        elif t == "[ Name ]":
            name_ph.append(shape)
        elif t == "[ Role ]":
            role_ph.append(shape)
        elif t.startswith("[ One line"):
            bio_ph.append(shape)
    team = [
        (
            "PO",
            "ProximaOpal",
            "Founder & Full-Stack Product Engineer",
            "Web developer and UI/UX designer based in Nairobi. Designed, built and shipped "
            "Aurora Travels solo: frontend, interactive maps, craft marketplace and live M-Pesa STK Push.",
        ),
        (
            "—",
            "Open role",
            "Tourism & hospitality partnerships lead",
            "Formalize shop, park and guide agreements across Kenya.",
        ),
        (
            "—",
            "Open role",
            "Artisan network coordinator",
            "Onboard and verify craft shops county by county.",
        ),
        (
            "—",
            "Open role",
            "Growth / GTM support",
            "Run the Nairobi + Maasai Mara pilot launch.",
        ),
    ]
    for i, (photo, name, role, bio) in enumerate(team):
        if i < len(photo_ph):
            set_text(photo_ph[i], photo)
        if i < len(name_ph):
            set_text(name_ph[i], name)
        if i < len(role_ph):
            set_text(role_ph[i], role)
        if i < len(bio_ph):
            set_text(bio_ph[i], bio)

    # ---- 11 Ask ----
    s = prs.slides[10]
    set_text(
        first_contains(s, "What you need"),
        "What Aurora Travels needs, and what it enables",
    )
    set_text(
        first_contains(s, "Funding / support requested"),
        "SEED / PARTNERSHIP CAPITAL\n"
        "US$180,000 ≈ KES 23.4 million (at ~130 KES/USD) — 18-month Nairobi + Maasai Mara pilot.\n"
        "Open to grant, equity, in-kind tourism partnerships, or blended support.\n"
        "USE OF FUNDS:\n"
        "• Product & engineering — 40% (US$72,000)\n"
        "• Artisan & guide network growth — 35% (US$63,000)\n"
        "• Marketing & pilot launch — 25% (US$45,000)",
    )
    set_text(
        first_contains(s, "12–24 month projection"),
        "12–24 MONTH PROJECTION\n"
        "Shop partners onboarded (model): 24 → 40 → 65 → 90 → 120 → 160 → 210 → 270 over 8 quarters.\n"
        "Y3 modeled platform revenue ~US$5.0M at 11% blended take on US$45M GMV.\n"
        "Headline metric: named Kenyan shops live on Aurora with M-Pesa checkout.\n"
        "Ask and growth curve are founder-proposed for the pilot — not achieved figures.",
    )


def main() -> None:
    if not SRC.exists():
        raise SystemExit(f"Missing template: {SRC}")

    # Keep a clean template copy in-repo
    TEMPLATE_COPY.write_bytes(SRC.read_bytes())

    prs = Presentation(str(SRC))
    assert len(prs.slides) == 11
    fill(prs)
    make_legible(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({OUT.stat().st_size} bytes)")

    # Verify
    prs2 = Presentation(str(OUT))
    leftovers = []
    for i, slide in enumerate(prs2.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = "\n".join(p.text for p in shape.text_frame.paragraphs)
            for needle in (
                "PITCH DECK TEMPLATE",
                "Replace the guidance",
                "[brief description]",
                "[ number ]",
                "[ Name ]",
                "[Insert",
                "What health challenge",
            ):
                if needle in t:
                    leftovers.append((i, needle))
    print("leftovers:", leftovers or "none")

    # Spot-check bold + dark on slide 2 body
    s2 = prs2.slides[1]
    bold_ok = dark_ok = 0
    runs = 0
    for shape in s2.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                runs += 1
                if run.font.bold:
                    bold_ok += 1
                try:
                    c = run.font.color.rgb
                    if c and c[0] < 80 and c[1] < 80 and c[2] < 80:
                        dark_ok += 1
                except Exception:
                    pass
    print(f"slide2 runs={runs} bold={bold_ok} dark={dark_ok}")


if __name__ == "__main__":
    main()
